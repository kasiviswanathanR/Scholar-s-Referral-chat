import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from dotenv import load_dotenv
import time
import docx
from pptx import Presentation


load_dotenv()
os.getenv("YOURS_API_KEY")
genai.configure(api_key=os.getenv("YOURS_API_KEY"))

def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_docx(docx_docs):
    text = ""
    for docx_file in docx_docs:
        doc = docx.Document(docx_file)
        for para in doc.paragraphs:
            text += para.text + '\n'
    return text

def extract_text_from_pptx(pptx_docs):
    text = ""
    for pptx_file in pptx_docs:
        prs = Presentation(pptx_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
    return text


def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks


def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")


def get_conversational_chain():
    prompt_template = """
    If the context not presented just use your ai model and answeer otherwise ,Answer the question as detailed as possible from the provided context process it and tell with your own words, make sure to provide all the details with accuracy and
    faster response with ur generative intelligence\n\n
    Context:\n {context}?\n
    Question: \n{question}\n

    Answer:
    """

    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)

    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)

    return chain


def user_input(user_question):
    chat_history = []

    # Load chat history
    chat_history = load_chat_history()

    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")

    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)

    chain = get_conversational_chain()

    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)

    chat_history.append(("You:", user_question))
    chat_history.append(("Gemini💁:", response["output_text"]))

    # Save the updated chat history
    save_chat_history(chat_history)

    return response["output_text"]



def save_chat_history(chat_history):
    with open("chat_history.txt", "w", encoding='utf-8') as file:
        for sender, message in chat_history:
            file.write(f"{sender} {message}\n")


def load_chat_history():
    chat_history = []
    if os.path.exists("chat_history.txt"):
        with open("chat_history.txt", "r", encoding="utf-8") as file:
            for line in file:
                line = line.strip()
                if line:
                    parts = line.split(" ", 1)
                    if len(parts) == 2:
                        sender, message = parts
                        chat_history.append((sender, message))
                    else:
                        # Handle case where there is no space-separated sender and message
                        # For example, if there's only a message without a sender
                        chat_history.append(("Unknown:", line))
    return chat_history

def empty_chat_history_file():
    with open("chat_history.txt", "w", encoding="utf-8") as file:
        file.truncate(0)


def main():
    st.set_page_config(
        page_title="Scholar's referral chat",
        page_icon="🤖"
    )

    # Set page title and initial layout
   
    chat_history = load_chat_history()

    # Sidebar for file upload and processing
    with st.sidebar:
        st.title("📤 Upload & 🤔 Inquire")
        uploaded_files = st.file_uploader("Upload your PDF Files and Click on the Submit & Process Button", type=["pdf","docx", "pptx"], accept_multiple_files=True)

        if st.button("Process Files"):
            with st.spinner("Processing..."):
                pdf_docs = [file for file in uploaded_files if file.type == "application/pdf"]
                pptx_docs = [file for file in uploaded_files if file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"]
                docx_docs = [file for file in uploaded_files if file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"]
                raw_text = ""

                if pdf_docs:
                    raw_text += get_pdf_text(pdf_docs)

                if docx_docs:
                    raw_text += extract_text_from_docx(docx_docs)

                if pptx_docs:
                    raw_text += extract_text_from_pptx(pptx_docs)

                text_chunks = get_text_chunks(raw_text)
                get_vector_store(text_chunks)
                st.session_state.last_action = "process"
                msg = st.empty()  # Create an empty placeholder

                msg.success("Done")

            # Sleep for 5 seconds
                time.sleep(5)
            
            # Clear the success message after 5 seconds
                msg.empty()

        
        # Call your function to clear text file here


    # Define CSS for the chat bubbles
    st.markdown(
        """
       <div style="display: flex; justify-content: center; margin-bottom: 20px;">
            <pre>
            <h1 style="color: #11111; text-align: center;padding-bottom-20px">Scholar's Referral-Chat</h1>
            <p>Dive Deeper into your documents! 📚🔍</p>
            <p>Explore the depths of knowledge and uncover hidden treasures. 🌊💎</p>
            </pre>
        </div>
        <style>
        .chat-container {
            display: flex;
            flex-direction: column;
            gap: 20px; /* Increased gap between messages */
            align-items: flex-start; /* Align bubbles to the left */
            width: 100%; /* Make the chat container full width */
        }
        .user-bubble {
            background-color: #f0f0f0;
            color: #000;
            padding: 20px;
            border-radius: 15px;
            max-width: 100%; /* Limit user bubble width */
            text-align: left; /* Align text to the left */
            margin-bottom: 10px;
            box-shadow: 2px 2px 5px rgba(0, 0, 0, 0.1); /* Add a slight shadow for depth */
        }
        .model-bubble {
            background-color: #0078d4;
            color: #fff;
            padding: 20px;
            border-radius: 15px;
            max-width: 100%; /* Limit model bubble width */
            text-align: left; /* Align text to the left */
            margin-bottom: 10px;
            box-shadow: 2px 2px 5px rgba(0, 0, 0, 0.1); /* Add a slight shadow for depth */
        }

        </style>
        """,
        unsafe_allow_html=True,
    )

    # Form for user input at the bottom 
    if user_question := st.chat_input("What is up?"):  
        response_text = user_input(user_question)
        chat_history.append(("👤 User:", user_question))
        chat_history.append(("🤖 AI:", response_text))
            # Save the updated chat history
        save_chat_history(chat_history)

    if st.sidebar.button('Clear Chat History'):
        empty_chat_history_file()

    # Display the conversation history
    for sender, message in chat_history:
        if sender == "👤 User:":
            st.markdown(f"<div class='user-bubble'>{sender} {message}</div>", unsafe_allow_html=True)
        else:
            st.markdown(f"<div class='model-bubble'>{sender} {message}</div>", unsafe_allow_html=True)


if __name__ == "__main__":
    main()